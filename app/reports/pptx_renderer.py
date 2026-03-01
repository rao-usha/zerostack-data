"""
PowerPoint (PPTX) Renderer for Nexdata Reports.

Builds an Investment Committee memo deck from report data.
Each section becomes 1 slide with a consistent navy-branded layout.
Charts use native PPTX chart objects. Tables use formatted PPTX tables.
"""

import logging
from io import BytesIO
from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
HEADER_H = Inches(0.85)
CONTENT_TOP = Inches(1.2)
CONTENT_W = Inches(12.133)  # SLIDE_W - 2*MARGIN

NAVY = RGBColor(0x1A, 0x36, 0x5D)
NAVY_LIGHT = RGBColor(0x2B, 0x6C, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50 = RGBColor(0xF7, 0xFA, 0xFC)
GRAY_100 = RGBColor(0xED, 0xF2, 0xF7)
GRAY_200 = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_500 = RGBColor(0x71, 0x80, 0x96)
GRAY_700 = RGBColor(0x4A, 0x55, 0x68)
GRAY_800 = RGBColor(0x2D, 0x37, 0x48)
ACCENT = RGBColor(0xED, 0x89, 0x36)
GREEN = RGBColor(0x38, 0xA1, 0x69)
RED = RGBColor(0xE5, 0x3E, 0x3E)
BLUE_LIGHT = RGBColor(0x63, 0xB3, 0xED)
PURPLE = RGBColor(0x80, 0x5A, 0xD5)
TEAL = RGBColor(0x31, 0x97, 0x95)

CHART_COLORS = [NAVY_LIGHT, BLUE_LIGHT, ACCENT, GREEN, GRAY_500, PURPLE, TEAL, RED]

FONT_NAME = "Calibri"
FONT_BODY = Pt(11)
FONT_SMALL = Pt(9)
FONT_TINY = Pt(8)


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def _fmt(n):
    if n is None:
        return "-"
    try:
        return f"{int(n):,}"
    except (ValueError, TypeError):
        return str(n)


def _fmt_currency(n):
    if n is None or n == 0:
        return "-"
    try:
        n = float(n)
    except (ValueError, TypeError):
        return str(n)
    if abs(n) >= 1_000_000_000:
        return f"${n / 1_000_000_000:.1f}B"
    if abs(n) >= 1_000_000:
        return f"${n / 1_000_000:.1f}M"
    if abs(n) >= 1_000:
        return f"${n / 1_000:.0f}K"
    return f"${n:,.0f}"


def _pct(n, decimals=0):
    if n is None:
        return "-"
    try:
        return f"{float(n):.{decimals}f}%"
    except (ValueError, TypeError):
        return str(n)


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def _set_text(shape, text, font_size=FONT_BODY, bold=False, color=GRAY_800,
              alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment


def _add_textbox(slide, left, top, width, height, text,
                 font_size=FONT_BODY, bold=False, color=GRAY_800,
                 alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_multiline(slide, left, top, width, height, lines,
                   font_size=FONT_BODY, color=GRAY_800, spacing=Pt(4)):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(line)
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing
    return txBox


def _add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Cover slide
# ---------------------------------------------------------------------------

def _add_cover(prs, title, subtitle, badge_text, date_str):
    """Full navy cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Navy background
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)

    # Thin accent line
    _add_rect(slide, Inches(4), Inches(2.6), Inches(5.333), Pt(3), ACCENT)

    # Title
    _add_textbox(slide, Inches(1), Inches(2.9), Inches(11.333), Inches(1.2),
                 title, font_size=Pt(36), bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_textbox(slide, Inches(1), Inches(4.1), Inches(11.333), Inches(0.6),
                 subtitle, font_size=Pt(16), color=RGBColor(0xA0, 0xAE, 0xC0),
                 alignment=PP_ALIGN.CENTER)

    # Badge
    if badge_text:
        _add_textbox(slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
                     badge_text, font_size=Pt(13), bold=True, color=ACCENT,
                     alignment=PP_ALIGN.CENTER)

    # Bottom bar
    _add_textbox(slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.5),
                 f"CONFIDENTIAL  ·  Nexdata Investment Intelligence  ·  {date_str}",
                 font_size=Pt(10), color=GRAY_500, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Section header bar (on every content slide)
# ---------------------------------------------------------------------------

def _section_slide(prs, number, title):
    """Create a content slide with the navy header bar. Returns the slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Navy header bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, NAVY)

    # Section number + title
    label = f"{number}. {title}" if number else title
    _add_textbox(slide, MARGIN, Inches(0.15), CONTENT_W, Inches(0.55),
                 label, font_size=Pt(20), bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)

    return slide


# ---------------------------------------------------------------------------
# KPI boxes
# ---------------------------------------------------------------------------

def _add_kpi_boxes(slide, metrics, top=CONTENT_TOP, left=MARGIN):
    """Add a row of KPI metric boxes.

    metrics: list of (label, value) or (label, value, detail) tuples
    """
    if not metrics:
        return top
    n = len(metrics)
    gap = Inches(0.15)
    total_gap = gap * (n - 1)
    box_w = int((CONTENT_W - total_gap) / n)
    box_h = Inches(0.9)

    for i, m in enumerate(metrics):
        label = m[0]
        value = str(m[1]) if len(m) > 1 else "-"
        x = left + i * (box_w + gap)

        # Box background
        rect = _add_rect(slide, x, top, box_w, box_h, GRAY_50)
        rect.line.color.rgb = GRAY_200
        rect.line.width = Pt(0.75)

        # Label
        _add_textbox(slide, x + Pt(8), top + Pt(6), box_w - Pt(16), Pt(14),
                     label.upper(), font_size=Pt(7.5), bold=True, color=GRAY_500,
                     alignment=PP_ALIGN.CENTER)

        # Value
        _add_textbox(slide, x + Pt(8), top + Pt(22), box_w - Pt(16), Pt(30),
                     value, font_size=Pt(22), bold=True, color=NAVY,
                     alignment=PP_ALIGN.CENTER)

    return top + box_h + Inches(0.2)


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------

def _add_table(slide, headers, rows, left=None, top=None, col_widths=None,
               max_rows=18):
    """Add a formatted table. Returns (table_shape, bottom_y)."""
    if left is None:
        left = MARGIN
    if top is None:
        top = CONTENT_TOP

    display_rows = rows[:max_rows]
    n_rows = len(display_rows) + 1  # +1 for header
    n_cols = len(headers)
    width = CONTENT_W
    row_h = Inches(0.28)
    tbl_h = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, tbl_h)
    table = shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    else:
        col_w = int(width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8.5)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME

    # Data rows
    for i, row in enumerate(display_rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if val is not None else "-"
            # Alternating row color
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_50
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8.5)
                paragraph.font.color.rgb = GRAY_800
                paragraph.font.name = FONT_NAME

    # Remove table borders for cleaner look
    for i in range(n_rows):
        for j in range(n_cols):
            cell = table.cell(i, j)
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return shape, top + tbl_h


# ---------------------------------------------------------------------------
# Native PPTX charts
# ---------------------------------------------------------------------------

def _add_bar_chart(slide, categories, values, title=None,
                   left=None, top=None, width=None, height=None,
                   horizontal=False, colors=None):
    """Add a native PPTX bar chart."""
    if left is None:
        left = MARGIN
    if top is None:
        top = CONTENT_TOP
    if width is None:
        width = CONTENT_W
    if height is None:
        height = Inches(4.5)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Value", values)

    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Style the chart
    plot = chart.plots[0]
    plot.gap_width = 80

    # Color the bars
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = NAVY_LIGHT

    # Axis styling
    if chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.category_axis.tick_labels.font.color.rgb = GRAY_700
        chart.category_axis.tick_labels.font.name = FONT_NAME
    if chart.value_axis:
        chart.value_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.color.rgb = GRAY_700
        chart.value_axis.tick_labels.font.name = FONT_NAME

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = GRAY_800

    return chart_frame


def _add_doughnut_chart(slide, categories, values,
                        left=None, top=None, width=None, height=None,
                        title=None, colors=None):
    """Add a native PPTX doughnut chart."""
    if left is None:
        left = MARGIN
    if top is None:
        top = CONTENT_TOP
    if width is None:
        width = Inches(5)
    if height is None:
        height = Inches(4)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Value", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = FONT_NAME
    chart.legend.include_in_layout = False

    # Color the segments
    if colors:
        plot = chart.plots[0]
        series = plot.series[0]
        for i, color in enumerate(colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    return chart_frame


# ---------------------------------------------------------------------------
# Callout box
# ---------------------------------------------------------------------------

def _add_callout(slide, text, top, left=None, width=None, variant="info"):
    """Add a callout/insight box."""
    if left is None:
        left = MARGIN
    if width is None:
        width = CONTENT_W
    h = Inches(0.6)

    accent = NAVY_LIGHT
    if variant == "warn":
        accent = ACCENT
    elif variant == "good":
        accent = GREEN

    # Accent bar
    _add_rect(slide, left, top, Pt(4), h, accent)

    # Background
    bg = _add_rect(slide, left + Pt(4), top, width - Pt(4), h, GRAY_50)
    bg.line.fill.background()

    # Text
    _add_textbox(slide, left + Pt(14), top + Pt(4), width - Pt(24), h - Pt(8),
                 text, font_size=Pt(9), color=GRAY_700)

    return top + h + Inches(0.1)


# ===========================================================================
# MedSpa Report Builder
# ===========================================================================

def _render_medspa(data: Dict[str, Any]) -> bytes:
    """Build the MedSpa Market IC memo deck."""
    from datetime import datetime

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    summary = data.get("summary", {})
    prospects_by_state = data.get("prospects_by_state", [])
    grade_dist = data.get("grade_distribution", [])
    top_targets = data.get("top_targets", [])
    zip_conc = data.get("zip_concentration", [])
    a_grade_state = data.get("a_grade_by_state", [])
    zip_affluence = data.get("zip_affluence_by_state", [])
    pe_comps = data.get("pe_comps", [])
    recent_deals = data.get("recent_deals", [])
    whitespace_by_state = data.get("whitespace_by_state", [])
    bls_wages = data.get("bls_wages", {})
    category_breakdown = data.get("category_breakdown", [])
    pe_financials = data.get("pe_financials", [])
    pe_fin_summary = data.get("pe_financial_summary", {})
    top_by_reviews = data.get("top_by_reviews", [])
    low_competition_gems = data.get("low_competition_gems", [])
    deal_model = data.get("deal_model", {})
    stealth_wealth = data.get("stealth_wealth", {})
    migration_alpha = data.get("migration_alpha", {})
    provider_density = data.get("provider_density", {})
    real_estate_alpha = data.get("real_estate_alpha", {})
    deposit_wealth = data.get("deposit_wealth", {})
    business_formation = data.get("business_formation", {})
    opportunity_zones = data.get("opportunity_zones", {})
    demographic_demand = data.get("demographic_demand", {})
    pe_competitive = data.get("pe_competitive", {})
    construction_momentum = data.get("construction_momentum", {})
    medical_cpi = data.get("medical_cpi", {})
    talent_pipeline = data.get("talent_pipeline", {})

    date_str = datetime.utcnow().strftime("%B %Y")

    # ── COVER ──────────────────────────────────────────────────────────
    _add_cover(
        prs,
        title="Aesthetics Roll-Up Thesis",
        subtitle="Med-Spa Market Analysis  ·  Acquisition Target Intelligence",
        badge_text=f"{_fmt(summary.get('total_prospects'))} Prospects Discovered",
        date_str=date_str,
    )

    # ── EXECUTIVE SUMMARY ──────────────────────────────────────────────
    slide = _section_slide(prs, 1, "Executive Summary")
    y = _add_kpi_boxes(slide, [
        ("Total Prospects", _fmt(summary.get("total_prospects"))),
        ("A-Grade Targets", _fmt(summary.get("a_grade"))),
        ("States Covered", _fmt(summary.get("states_covered"))),
        ("Avg Acq. Score", str(summary.get("avg_score", 0))),
        ("Avg Yelp Rating", f"{summary.get('avg_rating', 0):.1f}"),
    ])

    thesis_lines = [
        f"Nexdata identified {_fmt(summary.get('total_prospects'))} med-spa acquisition prospects across {_fmt(summary.get('states_covered'))} states.",
        f"{_fmt(summary.get('a_grade'))} A-grade and {_fmt(summary.get('ab_grade'))} A/B-grade targets scored via 5-factor composite model.",
        "Model inputs: ZIP affluence (IRS SOI), Yelp consumer signals, competitive density, price tier, review volume.",
    ]
    top5 = prospects_by_state[:5]
    if top5:
        thesis_lines.append(
            "Top markets: " + ", ".join(f"{s['state']} ({s['count']})" for s in top5)
        )
    _add_multiline(slide, MARGIN, y, CONTENT_W, Inches(2), thesis_lines,
                   font_size=Pt(12), color=GRAY_700)

    # Grade distribution as mini-table
    if grade_dist:
        y2 = y + Inches(1.8)
        gd_headers = ["Grade", "Count", "% of Total"]
        gd_rows = [[g["grade"], _fmt(g["count"]), f"{g['pct']}%"] for g in grade_dist]
        _add_table(slide, gd_headers, gd_rows, top=y2)

    # ── MARKET MAP ─────────────────────────────────────────────────────
    if prospects_by_state or grade_dist:
        slide = _section_slide(prs, 2, "Market Map")

        if prospects_by_state:
            cats = [s["state"] for s in prospects_by_state[:15]]
            vals = [float(s["count"]) for s in prospects_by_state[:15]]
            _add_bar_chart(slide, cats, vals,
                           title="Prospects by State (Top 15)",
                           left=MARGIN, top=CONTENT_TOP,
                           width=Inches(7.5), height=Inches(5.5),
                           horizontal=True)

        if grade_dist:
            g_labels = [f"Grade {g['grade']}" for g in grade_dist]
            g_values = [float(g["count"]) for g in grade_dist]
            g_colors = []
            cmap = {"A": NAVY_LIGHT, "B": GREEN, "C": ACCENT, "D": GRAY_500, "F": RED}
            for g in grade_dist:
                g_colors.append(cmap.get(g["grade"], GRAY_500))
            _add_doughnut_chart(slide, g_labels, g_values,
                                title="Grade Distribution",
                                left=Inches(8.3), top=CONTENT_TOP,
                                width=Inches(4.5), height=Inches(4),
                                colors=g_colors)

    # ── TOP ACQUISITION TARGETS ────────────────────────────────────────
    if top_targets:
        slide = _section_slide(prs, 3, "Top Acquisition Targets")
        _add_textbox(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(0.3),
                     f"{len(top_targets)} highest-scoring A/B-grade prospects ranked by composite score.",
                     font_size=FONT_SMALL, color=GRAY_500)

        headers = ["#", "Name", "City", "State", "Score", "Grade", "Rating", "Reviews"]
        rows = []
        for i, t in enumerate(top_targets[:20], 1):
            rows.append([
                str(i), t["name"], t.get("city", "-"), t.get("state", "-"),
                f"{t['score']:.0f}", t["grade"], f"{t['rating']:.1f}", _fmt(t["reviews"]),
            ])
        _add_table(slide, headers, rows, top=CONTENT_TOP + Inches(0.4), max_rows=20)

    # ── MARKET CONCENTRATION ───────────────────────────────────────────
    if zip_conc or a_grade_state:
        slide = _section_slide(prs, 4, "Market Concentration Analysis")
        y = CONTENT_TOP

        if zip_conc:
            _add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.25),
                         "Top ZIPs by Prospect Density", font_size=Pt(12),
                         bold=True, color=NAVY)
            y += Inches(0.3)
            headers = ["ZIP", "City", "State", "Prospects", "Avg Score", "ZIP Score"]
            rows = [[z["zip_code"], z.get("city", "-"), z.get("state", "-"),
                      str(z["prospect_count"]), str(z["avg_score"]),
                      str(z["avg_zip_score"])] for z in zip_conc[:12]]
            _, y = _add_table(slide, headers, rows, top=y, max_rows=12)

    # ── ZIP AFFLUENCE ──────────────────────────────────────────────────
    if zip_affluence:
        slide = _section_slide(prs, 5, "ZIP Affluence Profile")
        headers = ["State", "ZIPs Scored", "Avg Score", "Affluence", "Avg AGI", "A-Grade ZIPs"]
        rows = [[z["state"], str(z["zip_count"]), str(z["avg_score"]),
                  str(z["avg_affluence"]), _fmt_currency(z["avg_agi"]),
                  str(z["a_zips"])] for z in zip_affluence[:15]]
        _add_table(slide, headers, rows, top=CONTENT_TOP, max_rows=15)

    # ── COMPETITIVE LANDSCAPE ──────────────────────────────────────────
    if pe_comps or recent_deals:
        slide = _section_slide(prs, 6, "Competitive Landscape")
        y = CONTENT_TOP

        if pe_comps:
            _add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.25),
                         "PE-Backed Aesthetics Platforms", font_size=Pt(12),
                         bold=True, color=NAVY)
            y += Inches(0.3)
            headers = ["Company", "Industry", "PE Owner", "Employees", "Status"]
            rows = []
            for c in pe_comps[:12]:
                rows.append([
                    c["name"], c.get("industry", "-"),
                    c.get("pe_owner", "-"),
                    _fmt(c.get("employees")), c.get("status", "-"),
                ])
            _, y = _add_table(slide, headers, rows, top=y, max_rows=12)
            y += Inches(0.15)
            _add_callout(slide,
                         f"{len(pe_comps)} PE platforms identified — strong institutional interest in aesthetics.",
                         top=y)

    # ── METHODOLOGY ────────────────────────────────────────────────────
    slide = _section_slide(prs, 7, "Data Sources & Methodology")
    _add_textbox(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(0.3),
                 "Prospect Scoring Weights", font_size=Pt(12), bold=True, color=NAVY)
    p_weights = [
        ["ZIP Affluence (IRS SOI)", "30%"],
        ["Yelp Rating", "25%"],
        ["Review Volume", "20%"],
        ["Low Competition", "15%"],
        ["Price Tier", "10%"],
    ]
    _add_table(slide, ["Factor", "Weight"], p_weights,
               top=CONTENT_TOP + Inches(0.35), max_rows=10)

    _add_textbox(slide, MARGIN, Inches(4.0), CONTENT_W, Inches(0.3),
                 "ZIP Affluence Scoring Weights", font_size=Pt(12), bold=True, color=NAVY)
    z_weights = [
        ["Affluence Density", "30%"],
        ["Discretionary Wealth", "25%"],
        ["Market Size", "20%"],
        ["Professional Density", "15%"],
        ["Wealth Concentration", "10%"],
    ]
    _add_table(slide, ["Factor", "Weight"], z_weights,
               top=Inches(4.35), max_rows=10)

    # ── WHITESPACE ─────────────────────────────────────────────────────
    if whitespace_by_state:
        slide = _section_slide(prs, 8, "Whitespace Analysis")
        headers = ["State", "Underserved ZIPs", "Avg Affluence", "Opportunity Score"]
        rows = []
        for w in whitespace_by_state[:15]:
            rows.append([
                w.get("state", "-"), str(w.get("underserved_count", 0)),
                str(w.get("avg_affluence", 0)), str(w.get("opportunity_score", 0)),
            ])
        _add_table(slide, headers, rows, top=CONTENT_TOP, max_rows=15)

    # ── WORKFORCE ──────────────────────────────────────────────────────
    if bls_wages:
        slide = _section_slide(prs, 9, "Workforce Economics")
        wage_data = bls_wages.get("occupations", [])
        if wage_data:
            headers = ["Occupation", "Median Wage", "Employment", "Growth"]
            rows = [[w.get("title", "-"), _fmt_currency(w.get("median_wage")),
                      _fmt(w.get("employment")),
                      _pct(w.get("growth_pct"))] for w in wage_data[:12]]
            _add_table(slide, headers, rows, top=CONTENT_TOP, max_rows=12)

    # ── CATEGORIES ─────────────────────────────────────────────────────
    if category_breakdown:
        slide = _section_slide(prs, 10, "Service Category Breakdown")
        headers = ["Category", "Count", "% of Total", "Avg Score", "Avg Rating"]
        rows = [[c.get("category", "-"), _fmt(c.get("count")),
                  _pct(c.get("pct")), str(c.get("avg_score", "-")),
                  f"{c.get('avg_rating', 0):.1f}"] for c in category_breakdown[:15]]
        _add_table(slide, headers, rows, top=CONTENT_TOP, max_rows=15)

    # ── PE BENCHMARKS ──────────────────────────────────────────────────
    if pe_financials:
        slide = _section_slide(prs, 11, "PE Platform Benchmarking")
        if pe_fin_summary:
            y = _add_kpi_boxes(slide, [
                ("Avg Revenue", _fmt_currency(pe_fin_summary.get("avg_revenue"))),
                ("Avg EBITDA", _fmt_currency(pe_fin_summary.get("avg_ebitda"))),
                ("Avg Locations", _fmt(pe_fin_summary.get("avg_locations"))),
                ("Avg Margin", _pct(pe_fin_summary.get("avg_margin"), 1)),
            ])
        else:
            y = CONTENT_TOP

        headers = ["Platform", "Revenue", "EBITDA", "Locations", "Margin"]
        rows = [[p.get("name", "-"), _fmt_currency(p.get("revenue")),
                  _fmt_currency(p.get("ebitda")), _fmt(p.get("locations")),
                  _pct(p.get("margin"), 1)] for p in pe_financials[:10]]
        _add_table(slide, headers, rows, top=y, max_rows=10)

    # ── GROWTH SIGNALS ─────────────────────────────────────────────────
    if top_by_reviews or low_competition_gems:
        slide = _section_slide(prs, 12, "Review Velocity & Growth Signals")
        y = CONTENT_TOP

        if top_by_reviews:
            _add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.25),
                         "Highest Review Volume (Proven Demand)", font_size=Pt(12),
                         bold=True, color=NAVY)
            y += Inches(0.3)
            headers = ["Name", "City", "State", "Reviews", "Rating", "Grade"]
            rows = [[t["name"], t.get("city", "-"), t.get("state", "-"),
                      _fmt(t["reviews"]), f"{t['rating']:.1f}", t["grade"]]
                     for t in top_by_reviews[:8]]
            _, y = _add_table(slide, headers, rows, top=y, max_rows=8)

        if low_competition_gems:
            y += Inches(0.2)
            _add_textbox(slide, MARGIN, y, CONTENT_W, Inches(0.25),
                         "Low-Competition Gems (≤3 Competitors in ZIP)", font_size=Pt(12),
                         bold=True, color=NAVY)
            y += Inches(0.3)
            headers = ["Name", "City", "State", "Score", "Competitors"]
            rows = [[g["name"], g.get("city", "-"), g.get("state", "-"),
                      f"{g['score']:.0f}", str(g.get("competitors", "-"))]
                     for g in low_competition_gems[:8]]
            _add_table(slide, headers, rows, top=y, max_rows=8)

    # ── DEAL MODEL: UNIT ECONOMICS ─────────────────────────────────────
    dm_tiers = deal_model.get("tier_economics", [])
    if dm_tiers:
        slide = _section_slide(prs, 13, "Deal Model — Unit Economics")
        y = _add_kpi_boxes(slide, [
            ("A-Grade Locations", _fmt(deal_model.get("total_locations"))),
            ("Est. Revenue", _fmt_currency(deal_model.get("total_revenue"))),
            ("Est. EBITDA", _fmt_currency(deal_model.get("total_ebitda"))),
            ("Wtd. Margin", _pct(deal_model.get("weighted_margin", 0) * 100)),
        ])
        headers = ["Tier", "Locs", "Avg Rev", "Margin", "Avg EBITDA", "Entry Mult.", "Total Rev", "Total EBITDA"]
        rows = [[t["tier"], _fmt(t["count"]), _fmt_currency(t["avg_revenue"]),
                  f"{t['ebitda_margin']:.0%}", _fmt_currency(t["avg_ebitda"]),
                  f"{t['entry_multiple']:.1f}x", _fmt_currency(t["total_revenue"]),
                  _fmt_currency(t["total_ebitda"])] for t in dm_tiers]
        _add_table(slide, headers, rows, top=y, max_rows=10)

    # ── DEAL MODEL: CAPITAL ────────────────────────────────────────────
    dm_cap = deal_model.get("capital_stack", {})
    if dm_cap:
        slide = _section_slide(prs, 14, "Deal Model — Capital Requirements")
        y = _add_kpi_boxes(slide, [
            ("Total Acq. Cost", _fmt_currency(deal_model.get("total_acquisition_cost"))),
            ("Equity (40%)", _fmt_currency(dm_cap.get("equity"))),
            ("Debt (60%)", _fmt_currency(dm_cap.get("debt"))),
            ("Txn Costs (5%)", _fmt_currency(dm_cap.get("transaction_costs"))),
            ("Total Capital", _fmt_currency(dm_cap.get("total_capital_required"))),
        ])
        if dm_tiers:
            headers = ["Tier", "Locations", "Avg EBITDA", "Entry Mult.", "Total Acq. Cost"]
            rows = [[t["tier"], _fmt(t["count"]), _fmt_currency(t["avg_ebitda"]),
                      f"{t['entry_multiple']:.1f}x",
                      _fmt_currency(t.get("total_acq_cost"))] for t in dm_tiers]
            _add_table(slide, headers, rows, top=y, max_rows=10)

    # ── DEAL MODEL: RETURNS ────────────────────────────────────────────
    dm_scenarios = deal_model.get("scenarios", {})
    if dm_scenarios:
        slide = _section_slide(prs, 15, "Deal Model — Returns Analysis")
        y = CONTENT_TOP
        headers = ["Metric", "Conservative", "Base Case", "Aggressive"]
        metrics = ["entry_multiple", "exit_multiple", "hold_years",
                    "exit_revenue", "exit_ebitda", "equity_moic", "irr"]
        labels = ["Entry Multiple", "Exit Multiple", "Hold Period (Yrs)",
                   "Exit Revenue", "Exit EBITDA", "Equity MOIC", "IRR"]

        rows = []
        for label, key in zip(labels, metrics):
            row = [label]
            for sc in ["conservative", "base", "aggressive"]:
                s = dm_scenarios.get(sc, {})
                v = s.get(key, "-")
                if key in ("exit_revenue", "exit_ebitda"):
                    row.append(_fmt_currency(v) if v != "-" else "-")
                elif key == "irr":
                    row.append(f"{v:.0%}" if isinstance(v, (int, float)) else str(v))
                elif key == "equity_moic":
                    row.append(f"{v:.1f}x" if isinstance(v, (int, float)) else str(v))
                elif key in ("entry_multiple", "exit_multiple"):
                    row.append(f"{v:.1f}x" if isinstance(v, (int, float)) else str(v))
                else:
                    row.append(str(v))
            rows.append(row)
        _add_table(slide, headers, rows, top=y, max_rows=10)

    # ── SIGNAL SECTIONS (16-27) ────────────────────────────────────────
    # These follow a common pattern: KPIs + state rankings table

    _signal_section(prs, 16, "Stealth Wealth Signal", stealth_wealth, [
        ("ZIPs Analyzed", "summary.total_returns"),
        ("Stealth ZIPs", "summary.total_stealth"),
        ("Validated", "summary.validated"),
        ("Avg Non-Wage", "summary.avg_non_wage_income", "currency"),
    ], table_key="stealth_zips", table_headers=["ZIP", "State", "Non-Wage/Ret", "Non-Wage %", "Avg AGI", "Score"],
       table_fields=["zip_code", "state", "non_wage_per_return", "non_wage_pct", "avg_agi", "medspa_score"],
       table_fmt={"non_wage_per_return": "currency", "avg_agi": "currency", "non_wage_pct": "pct"})

    _signal_section(prs, 17, "Migration Alpha", migration_alpha, [
        ("Counties Tracked", "summary.counties_tracked"),
        ("Net Positive", "summary.net_positive"),
        ("Avg Net Flow", "summary.avg_net_flow", "currency"),
        ("Top Inflow State", "summary.top_inflow_state"),
    ], table_key="top_inflows", table_headers=["State", "County", "Net Flow", "Inflows", "Outflows"],
       table_fields=["state", "county", "net_flow", "inflows", "outflows"],
       table_fmt={"net_flow": "currency", "inflows": "currency", "outflows": "currency"})

    _signal_section(prs, 18, "Medical Provider Density", provider_density, [
        ("Providers Tracked", "summary.total_providers"),
        ("Opportunity ZIPs", "summary.opportunity_zips"),
        ("Avg Provider/ZIP", "summary.avg_per_zip"),
        ("Beneficiaries", "summary.total_beneficiaries"),
    ], table_key="opportunity_zips", table_headers=["ZIP", "State", "Providers", "Beneficiaries", "Score"],
       table_fields=["zip_code", "state", "provider_count", "beneficiaries", "medspa_score"])

    _signal_section(prs, 19, "Real Estate Appreciation Alpha", real_estate_alpha, [
        ("ZIPs Tracked", "summary.zips_tracked"),
        ("Avg Appreciation", "summary.avg_appreciation", "pct"),
        ("Top State", "summary.top_state"),
    ], table_key="top_zips", table_headers=["ZIP", "State", "Median Price", "YoY Change", "Score"],
       table_fields=["zip_code", "state", "median_price", "yoy_change", "medspa_score"],
       table_fmt={"median_price": "currency", "yoy_change": "pct"})

    _signal_section(prs, 20, "Deposit Wealth Concentration", deposit_wealth, [
        ("Branches Tracked", "summary.branches_tracked"),
        ("Total Deposits", "summary.total_deposits", "currency"),
        ("Avg Deposit/Cap", "summary.avg_per_capita", "currency"),
    ], table_key="top_states", table_headers=["State", "Deposits", "Per Capita", "Branches"],
       table_fields=["state", "total_deposits", "per_capita", "branches"],
       table_fmt={"total_deposits": "currency", "per_capita": "currency"})

    _signal_section(prs, 21, "Business Formation Velocity", business_formation, [
        ("ZIPs Tracked", "summary.zips_tracked"),
        ("Avg Biz Income", "summary.avg_biz_income", "currency"),
        ("Top State", "summary.top_state"),
    ], table_key="top_states", table_headers=["State", "Biz Income", "Density", "Growth"],
       table_fields=["state", "biz_income", "density", "growth"],
       table_fmt={"biz_income": "currency", "growth": "pct"})

    _signal_section(prs, 22, "Opportunity Zone Overlay", opportunity_zones, [
        ("OZ Tracts", "summary.total_tracts"),
        ("OZ with Medspas", "summary.tracts_with_medspas"),
        ("Avg Score in OZ", "summary.avg_score"),
    ], table_key="top_tracts", table_headers=["Tract", "State", "Designation", "Medspas", "Score"],
       table_fields=["tract_id", "state", "designation_type", "medspa_count", "avg_score"])

    _signal_section(prs, 23, "Demographic Demand Model", demographic_demand, [
        ("States Analyzed", "summary.states_analyzed"),
        ("Target Demo Pop", "summary.target_population"),
        ("Underserved States", "summary.underserved_count"),
    ], table_key="state_rankings", table_headers=["State", "Target Pop", "Medspas", "Per Capita", "Gap Score"],
       table_fields=["state", "target_pop", "medspa_count", "per_capita", "gap_score"],
       table_fmt={"target_pop": "number"})

    _signal_section(prs, 24, "PE Competitive Heat Map", pe_competitive, [
        ("Total Deals", "summary.total_deals"),
        ("Active Platforms", "summary.active_platforms"),
        ("Avg Deal Size", "summary.avg_deal_size", "currency"),
    ], table_key="recent_activity", table_headers=["Platform", "Target", "Year", "Type", "Size"],
       table_fields=["platform", "target", "year", "deal_type", "deal_size"],
       table_fmt={"deal_size": "currency"})

    _signal_section(prs, 25, "Construction Momentum", construction_momentum, [
        ("States Tracked", "summary.states_tracked"),
        ("Avg Permit Growth", "summary.avg_growth", "pct"),
        ("Top Growth State", "summary.top_state"),
    ], table_key="state_rankings", table_headers=["State", "Permits", "YoY Growth", "5yr CAGR"],
       table_fields=["state", "permits", "yoy_growth", "cagr_5yr"],
       table_fmt={"yoy_growth": "pct", "cagr_5yr": "pct"})

    _signal_section(prs, 26, "Medical CPI Pricing Power", medical_cpi, [
        ("Medical CPI", "summary.medical_cpi"),
        ("General CPI", "summary.general_cpi"),
        ("Spread", "summary.spread", "pct"),
    ], table_key="trend_data", table_headers=["Year", "Medical CPI", "General CPI", "Spread"],
       table_fields=["year", "medical_cpi", "general_cpi", "spread"],
       table_fmt={"spread": "pct"})

    _signal_section(prs, 27, "Talent Pipeline Pressure", talent_pipeline, [
        ("Openings", "summary.total_openings"),
        ("Hires", "summary.total_hires"),
        ("Ratio", "summary.openings_to_hires"),
        ("Trend", "summary.trend"),
    ], table_key="occupation_data", table_headers=["Occupation", "Openings", "Hires", "Ratio", "Quits Rate"],
       table_fields=["occupation", "openings", "hires", "ratio", "quits_rate"],
       table_fmt={"quits_rate": "pct"})

    # ── DISCLAIMER SLIDE ───────────────────────────────────────────────
    slide = _section_slide(prs, None, "Disclaimers & Data Sources")
    lines = [
        "This report does not constitute investment advice. All data sourced from public APIs.",
        "Prospect data: Yelp Business Search API. ZIP affluence: IRS SOI ZIP-level data.",
        "PE landscape: SEC EDGAR filings, company websites. Deal model uses AmSpa/IBISWorld benchmarks.",
        "Stealth Wealth: IRS SOI non-wage income. Migration: IRS county flows. Provider: CMS Medicare.",
        "Real Estate: Redfin/FHFA. Deposits: FDIC. Construction: HUD SOCDS. CPI: BLS.",
        "",
        f"Generated by Nexdata Investment Intelligence · {date_str}",
    ]
    _add_multiline(slide, MARGIN, CONTENT_TOP, CONTENT_W, Inches(5),
                   lines, font_size=Pt(11), color=GRAY_500, spacing=Pt(8))

    # ── SAVE ───────────────────────────────────────────────────────────
    buf = BytesIO()
    prs.save(buf)
    pdf_bytes = buf.getvalue()
    logger.info(f"PPTX generated: {len(pdf_bytes):,} bytes")
    return pdf_bytes


def _signal_section(prs, number, title, data, kpi_defs, table_key=None,
                    table_headers=None, table_fields=None, table_fmt=None):
    """Generic builder for signal sections (16-27).

    kpi_defs: list of (label, dotted_path) or (label, dotted_path, fmt_type)
    """
    if not data:
        return  # skip empty sections

    slide = _section_slide(prs, number, title)

    # KPI boxes
    metrics = []
    for kdef in kpi_defs:
        label = kdef[0]
        path = kdef[1]
        fmt_type = kdef[2] if len(kdef) > 2 else None

        # Navigate dotted path
        val = data
        for part in path.split("."):
            if isinstance(val, dict):
                val = val.get(part)
            else:
                val = None
                break

        if fmt_type == "currency":
            display = _fmt_currency(val)
        elif fmt_type == "pct":
            display = _pct(val, 1)
        elif fmt_type == "number":
            display = _fmt(val)
        else:
            display = _fmt(val) if isinstance(val, (int, float)) else (str(val) if val else "-")

        metrics.append((label, display))

    y = _add_kpi_boxes(slide, metrics)

    # Table
    if table_key and table_headers and table_fields:
        rows_data = data.get(table_key, [])
        if isinstance(rows_data, list) and rows_data:
            fmt_map = table_fmt or {}
            rows = []
            for item in rows_data[:15]:
                row = []
                for field in table_fields:
                    v = item.get(field) if isinstance(item, dict) else None
                    f = fmt_map.get(field)
                    if f == "currency":
                        row.append(_fmt_currency(v))
                    elif f == "pct":
                        row.append(_pct(v, 1))
                    elif f == "number":
                        row.append(_fmt(v))
                    else:
                        row.append(str(v) if v is not None else "-")
                rows.append(row)
            _add_table(slide, table_headers, rows, top=y, max_rows=15)


# ===========================================================================
# Public API
# ===========================================================================

def render_pptx(template_name: str, data: Dict[str, Any]) -> bytes:
    """Render report data as a PowerPoint deck.

    Currently supports: medspa_market
    """
    if template_name == "medspa_market":
        return _render_medspa(data)
    raise ValueError(
        f"PPTX export not yet supported for template '{template_name}'. "
        f"Supported: medspa_market"
    )
